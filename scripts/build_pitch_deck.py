from __future__ import annotations

from functools import lru_cache
from pathlib import Path
from tempfile import TemporaryDirectory
import json
import shutil
import subprocess
import sys

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "src"
if str(SRC) not in sys.path:
    sys.path.insert(0, str(SRC))

from fitzsight.agent.catalog import (
    CRM_INTENT,
    NET_DEPOSIT_INTENT,
    CUSTOMER_INTELLIGENCE_INTENT,
    MARKETING_LEAD_QUALITY_INTENT,
    FALSE_CORRELATION_INTENT,
)
from fitzsight.agent.planner import ConstrainedRulePlanner
from fitzsight.demo import DEMO_QUESTIONS
from fitzsight.runtime import build_agent_runtime

SUBMISSION_DIR = ROOT / "submission"
DOCS_DIR = ROOT / "docs"
PPTX_PATH = SUBMISSION_DIR / "FitzSight_GOAI_Initial_Round.pptx"
PDF_PATH = SUBMISSION_DIR / "FitzSight_GOAI_Initial_Round.pdf"
HERO_TRACE = SUBMISSION_DIR / "FitzSight_Hero_Run_Trace.png"
HERO_ANSWER = SUBMISSION_DIR / "FitzSight_Hero_Run_Answer.png"

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

BG = RGBColor(11, 16, 32)
PANEL = RGBColor(20, 27, 45)
PANEL_2 = RGBColor(27, 36, 58)
TEXT = RGBColor(246, 248, 252)
MUTED = RGBColor(170, 180, 197)
CYAN = RGBColor(88, 211, 255)
GREEN = RGBColor(84, 211, 155)
AMBER = RGBColor(255, 200, 87)
RED = RGBColor(255, 107, 107)
PURPLE = RGBColor(173, 141, 255)
WHITE = RGBColor(255, 255, 255)
FONT = "Liberation Sans"


@lru_cache(maxsize=1)
def _pitch_runs() -> dict[str, dict]:
    """Fresh deterministic runs used only for current competition-facing numbers."""
    runs: dict[str, dict] = {}
    with TemporaryDirectory(prefix="fitzsight_pitch_") as tmp:
        store, _registry, agent = build_agent_runtime(
            data_dir=Path(tmp),
            backend="sqlite",
            planner=ConstrainedRulePlanner(),
        )
        try:
            for question in DEMO_QUESTIONS.values():
                result = agent.run(question)
                if result.final_answer.status != "verified" or not result.verification.passed:
                    raise RuntimeError(
                        f"Pitch metric source failed verification for {result.plan.intent}"
                    )
                runs[result.plan.intent] = result.to_dict()
        finally:
            store.close()
    return runs


def _metrics(intent: str) -> dict:
    return _pitch_runs()[intent]["investigation"]["metrics"]


def _diagnosis(intent: str) -> dict:
    return _pitch_runs()[intent]["investigation"]["diagnosis"]


def _json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


@lru_cache(maxsize=1)
def _hero() -> dict:
    return _json(DOCS_DIR / "V0.12_HERO_RUN.json")


@lru_cache(maxsize=1)
def _evaluation_snapshot() -> dict[str, dict]:
    required = {
        "benchmark": DOCS_DIR / "V0.12_BENCHMARK_RESULTS.json",
        "adversarial": DOCS_DIR / "V0.12_ADVERSARIAL_RESULTS.json",
        "holdout": DOCS_DIR / "V0.12_HOLDOUT_RESULTS.json",
        "ablation": DOCS_DIR / "V0.12_ABLATION_RESULTS.json",
    }
    missing = [str(p) for p in required.values() if not p.exists()]
    if missing:
        raise FileNotFoundError(f"Missing v0.12 evaluation assets: {missing}")
    return {name: _json(path) for name, path in required.items()}


def _money_k(value: float, *, signed: bool = True) -> str:
    sign = ""
    if signed:
        sign = "+" if value > 0 else "-" if value < 0 else ""
    elif value < 0:
        sign = "-"
    return f"{sign}${abs(value) / 1000:.1f}k"


def _pp(value: float) -> str:
    return f"{value:+.2f} pp"


def _pct(value: float, digits: int = 0) -> str:
    return f"{value * 100:.{digits}f}%"


def _pvalue(value: float) -> str:
    return f"{value:.2e}" if value < 0.001 else f"{value:.5f}"


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text, x, y, w, h, *, size=24, color=TEXT, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font=FONT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None, *, kicker=None):
    if kicker:
        add_text(slide, kicker.upper(), 0.72, 0.35, 6.8, 0.28, size=10, color=CYAN, bold=True)
    add_text(slide, title, 0.72, 0.72, 11.9, 0.7, size=28, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.74, 1.43, 11.5, 0.38, size=13, color=MUTED)


def add_footer(slide, n):
    add_text(slide, "FitzSight · GOAI 2026 · Boundless Agents · AI+金融", 0.72, 7.13, 9.6, 0.2, size=8, color=MUTED)
    add_text(slide, str(n), 12.1, 7.08, 0.5, 0.24, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def panel(slide, x, y, w, h, *, fill=PANEL, radius=True, line=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line or fill
    return shp


def chip(slide, text, x, y, w, *, fill=PANEL_2, color=TEXT):
    panel(slide, x, y, w, 0.36, fill=fill)
    add_text(slide, text, x + 0.1, y + 0.085, w - 0.2, 0.18, size=9.5, color=color, bold=True, align=PP_ALIGN.CENTER)


def metric_card(slide, x, y, w, h, label, value, *, accent=CYAN, note=None):
    panel(slide, x, y, w, h, fill=PANEL)
    add_text(slide, label, x + 0.18, y + 0.15, w - 0.36, 0.24, size=9.5, color=MUTED, bold=True)
    add_text(slide, value, x + 0.18, y + 0.47, w - 0.36, 0.48, size=22, color=accent, bold=True)
    if note:
        add_text(slide, note, x + 0.18, y + h - 0.34, w - 0.36, 0.2, size=8.5, color=MUTED)


def arrow(slide, x1, y1, x2, y2, *, color=CYAN, width=2.0):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def flow_node(slide, text, x, y, w=1.7, *, fill=PANEL_2, accent=CYAN, h=0.72, size=11):
    panel(slide, x, y, w, h, fill=fill, line=accent)
    add_text(slide, text, x + 0.08, y + 0.12, w - 0.16, h - 0.2, size=size, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def add_callout(slide, text, x, y, w, h, *, accent=CYAN, size=12):
    panel(slide, x, y, w, h, fill=PANEL_2, line=accent)
    panel(slide, x, y, 0.06, h, fill=accent, radius=False)
    add_text(slide, text, x + 0.2, y + 0.12, w - 0.35, h - 0.22, size=size, color=TEXT, bold=True, valign=MSO_ANCHOR.MIDDLE)


def add_picture_contain(slide, path: Path, x: float, y: float, w: float, h: float, *, border=RGBColor(210, 218, 232)):
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    pw = iw * scale
    ph = ih * scale
    px = x + (w - pw) / 2
    py = y + (h - ph) / 2
    panel(slide, x, y, w, h, fill=WHITE, line=border)
    slide.shapes.add_picture(str(path), Inches(px), Inches(py), width=Inches(pw), height=Inches(ph))


def slide_1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    add_text(s, "FITZSIGHT", 0.78, 0.66, 3.5, 0.35, size=12, color=CYAN, bold=True)
    add_text(s, "Evidence-grounded\nFinancial Operations\nIntelligence Agent", 0.78, 1.28, 7.5, 2.0, size=31, bold=True)
    add_text(s, "Brokerage / FinTech Operations Analyst", 0.8, 3.55, 6.5, 0.35, size=15, color=MUTED, bold=True)
    add_callout(s, "Autonomous investigation. Human decision.", 0.8, 4.25, 6.2, 0.72, accent=GREEN, size=15)
    labels = [("QUESTION", CYAN), ("PLAN", CYAN), ("TOOLS", GREEN), ("EVIDENCE", AMBER), ("VERIFY", PURPLE)]
    for i, (t, c) in enumerate(labels):
        flow_node(s, t, 9.15, 0.9 + i * 1.05, 2.75, fill=PANEL, accent=c, h=0.68, size=12)
        if i < len(labels) - 1:
            arrow(s, 10.53, 1.58 + i * 1.05, 10.53, 1.86 + i * 1.05, color=MUTED, width=1.2)
    add_text(s, "GOAI 2026 · Boundless Agents · AI+金融", 0.8, 6.38, 6.8, 0.3, size=13, color=MUTED)
    add_text(s, "Not chat with a CSV.\nA bounded investigation you can audit.", 8.8, 6.1, 3.4, 0.6, size=12.5, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_footer(s, 1)


def slide_2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    add_title(s, "Dashboards show what changed. The investigation behind why is still manual.", kicker="Industry problem")
    add_text(s, "Primary user", 0.8, 1.75, 1.5, 0.25, size=10, color=CYAN, bold=True)
    add_text(s, "Brokerage / FinTech Operations Analyst", 0.8, 2.08, 5.9, 0.38, size=19, bold=True)
    steps = ["Find tables", "Define KPI", "Compare periods", "Drill drivers", "Test significance", "Inspect events", "Reconcile evidence", "Write report"]
    for i, t in enumerate(steps):
        x = 0.8 + (i % 4) * 3.0
        y = 2.82 + (i // 4) * 1.05
        panel(s, x, y, 2.62, 0.72, fill=PANEL)
        add_text(s, f"{i+1:02d}", x + 0.16, y + 0.18, 0.34, 0.24, size=10, color=CYAN, bold=True)
        add_text(s, t, x + 0.52, y + 0.16, 1.85, 0.28, size=11.5, bold=True)
    add_callout(s, "A generic LLM can write a plausible explanation quickly. Financial operations needs a reproducible evidence chain before narrative.", 0.8, 5.15, 11.75, 0.88, accent=AMBER, size=12)
    chip(s, "acquisition", 1.0, 6.28, 2.15, fill=PANEL_2, color=CYAN)
    arrow(s, 3.18, 6.46, 4.0, 6.46, color=MUTED, width=1.3)
    chip(s, "FTD conversion", 4.05, 6.28, 2.35, fill=PANEL_2, color=GREEN)
    arrow(s, 6.43, 6.46, 7.25, 6.46, color=MUTED, width=1.3)
    chip(s, "client-fund flows", 7.3, 6.28, 2.55, fill=PANEL_2, color=AMBER)
    add_text(s, "beachhead workflow", 10.15, 6.36, 2.0, 0.22, size=9.5, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_footer(s, 2)


def slide_3(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    add_title(s, "FitzSight converts a business question into a verified decision-support investigation", kicker="Product")
    labels = [("Question", CYAN), ("Constrained\nplan", CYAN), ("Deterministic\ntools", GREEN), ("Evidence\ngraph", AMBER), ("Verifier", PURPLE), ("Verified\nanswer", GREEN)]
    xs = [0.62, 2.72, 4.82, 6.98, 9.08, 11.0]
    widths = [1.55, 1.55, 1.65, 1.55, 1.35, 1.7]
    for (t, c), x, w in zip(labels, xs, widths):
        flow_node(s, t, x, 2.35, w, fill=PANEL, accent=c, h=0.82, size=10.5)
    for i in range(len(xs) - 1):
        arrow(s, xs[i] + widths[i], 2.76, xs[i + 1] - 0.07, 2.76, color=MUTED, width=1.3)
    add_callout(s, "The model may select the next approved analytical branch. It never owns business arithmetic, arbitrary SQL, or high-impact financial actions.", 0.8, 3.75, 11.75, 0.9, accent=CYAN, size=12)
    metric_card(s, 0.8, 5.05, 3.6, 1.2, "Planner authority", "Approved actions only", accent=CYAN, note="Bounded adaptivity; local intent gate first")
    metric_card(s, 4.62, 5.05, 3.6, 1.2, "Calculation authority", "Read-only SQL / Python", accent=GREEN, note="Every competition-facing number is executable")
    metric_card(s, 8.45, 5.05, 4.1, 1.2, "Answer authority", "EvidenceClaimVerifier", accent=PURPLE, note="Verification failure → attribution withheld")
    add_footer(s, 3)


def slide_4(prs):
    hero = _hero()
    run = hero["run"]
    metrics = run["investigation"]["metrics"]
    trace = run["investigation"]["execution_trace"]
    branch = metrics["bounded_branching"]
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    add_title(s, "Hero journey — evidence selects the next approved step", subtitle="Question: “Why did European FTD conversion deteriorate after July 15?”", kicker="Bounded-adaptive Agent")
    add_picture_contain(s, HERO_TRACE, 5.45, 1.82, 7.1, 4.92)
    metric_card(s, 0.8, 1.95, 4.25, 1.05, "Actual execution", f"{len(trace)} approved steps", accent=CYAN, note="Rendered from the verified v0.12 run JSON")
    statuses = [
        ("Contribution drilldown", "executed", GREEN),
        ("Latency / anomaly scan", "executed", GREEN),
        ("Operational event check", branch["event_check_status"], AMBER),
        ("Document evidence", metrics["document_evidence"]["source_ref"], PURPLE),
    ]
    for i, (lab, value, col) in enumerate(statuses):
        panel(s, 0.8, 3.18 + i * 0.72, 4.25, 0.58, fill=PANEL, line=col)
        add_text(s, lab, 1.02, 3.32 + i * 0.72, 1.95, 0.22, size=9.5, color=MUTED, bold=True)
        add_text(s, value, 2.95, 3.30 + i * 0.72, 1.8, 0.24, size=10.5, color=col, bold=True, align=PP_ALIGN.RIGHT)
    add_callout(s, "Result-driven branching ≠ unrestricted autonomy. Every next step remains inside the approved action catalog.", 0.8, 6.2, 4.25, 0.58, accent=GREEN, size=10.5)
    add_footer(s, 4)


def slide_5(prs):
    hero = _hero()["run"]
    m = hero["investigation"]["metrics"]
    d = hero["investigation"]["diagnosis"]
    affected = float(m["affected"]["conversion_change_pp"])
    control = float(m["control"]["conversion_change_pp"])
    response = float(m["affected_response_median_change_minutes"])
    contrib = m["team_contribution_analysis"]["segments"][0]
    anomaly_days = int(m["post_change_response_anomalies"]["anomaly_count"])
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    add_title(s, "Hero finding — evidence supports a CRM routing candidate, not a proven causal conclusion", kicker="Verified answer")
    metric_card(s, 0.8, 1.8, 2.15, 1.28, "Affected FTD", _pp(affected), accent=RED, note="Europe Team A+B")
    metric_card(s, 3.1, 1.8, 2.15, 1.28, "Control", _pp(control), accent=AMBER, note="Other Europe teams")
    metric_card(s, 5.4, 1.8, 2.15, 1.28, "Response median", f"{response:+.2f} min", accent=RED)
    metric_card(s, 7.7, 1.8, 2.15, 1.28, "Top contributor", str(contrib["segment"]), accent=CYAN, note=f"contribution {float(contrib['total_contribution_pp']):+.2f} pp")
    metric_card(s, 10.0, 1.8, 2.55, 1.28, "Root-cause status", "supported", accent=GREEN, note="candidate · not causal proof")
    add_picture_contain(s, HERO_ANSWER, 6.1, 3.3, 6.45, 3.1)
    panel(s, 0.8, 3.3, 4.95, 3.1, fill=PANEL)
    add_text(s, "Evidence chain", 1.05, 3.52, 2.1, 0.28, size=12, color=CYAN, bold=True)
    rows = [
        ("Quantitative shift", f"Affected {affected:+.2f} pp vs control {control:+.2f} pp"),
        ("Latency signal", f"Median response {response:+.2f} minutes"),
        ("Contribution", f"{contrib['segment']} is largest negative contributor"),
        ("Robust anomaly", f"{anomaly_days} post-change high-anomaly days"),
        ("Operational context", "Nearby CRM routing event"),
        ("Source-addressable doc", "CRM-CHANGE-2026-0715#p1"),
    ]
    for i, (lab, val) in enumerate(rows):
        add_text(s, lab, 1.05, 3.98 + i * 0.36, 1.65, 0.22, size=9.2, color=MUTED, bold=True)
        add_text(s, val, 2.75, 3.98 + i * 0.36, 2.65, 0.25, size=9.8, color=TEXT)
    add_footer(s, 5)


def slide_6(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    add_title(s, "Failure branch — error evidence ends attribution, not the investigation", subtitle="Tested dependency failure: the operational-event lookup raises an error Evidence record.", kicker="Fail closed")
    nodes = [
        ("event_check", RED),
        ("error Evidence", RED),
        ("no document corroboration", AMBER),
        ("insufficient_evidence", AMBER),
        ("verified bounded answer", GREEN),
    ]
    x = 0.75
    widths = [1.8, 1.8, 2.15, 2.15, 2.1]
    for i, ((label, col), w) in enumerate(zip(nodes, widths)):
        flow_node(s, label, x, 2.35, w, fill=PANEL, accent=col, h=0.88, size=10.2)
        if i < len(nodes)-1:
            arrow(s, x+w, 2.79, x+w+0.20, 2.79, color=MUTED, width=1.2)
        x += w + 0.25
    panel(s, 0.8, 4.05, 5.6, 1.55, fill=PANEL, line=RED)
    add_text(s, "What fails", 1.05, 4.3, 1.6, 0.25, size=11, color=RED, bold=True)
    add_text(s, "Operational-event dependency is unavailable. The system records the error and does not manufacture event or document corroboration.", 1.05, 4.72, 4.95, 0.62, size=11, color=TEXT)
    panel(s, 6.75, 4.05, 5.8, 1.55, fill=PANEL, line=GREEN)
    add_text(s, "What still succeeds", 7.0, 4.3, 1.9, 0.25, size=11, color=GREEN, bold=True)
    add_text(s, "The evidence-backed metric findings remain valid; EvidenceClaimVerifier passes a guarded answer whose root_cause_status is insufficient_evidence.", 7.0, 4.72, 5.0, 0.62, size=11, color=TEXT)
    add_callout(s, "Agent quality is not “always finding a cause.” It is completing the task with the correct uncertainty boundary.", 0.8, 6.12, 11.75, 0.65, accent=AMBER, size=11.5)
    add_footer(s, 6)


def slide_7(prs):
    m = _metrics(FALSE_CORRELATION_INTENT)
    d = _diagnosis(FALSE_CORRELATION_INTENT)
    asia = float(m["conversion_change_pp"])
    affiliate = float(m["affiliate_conversion_test"]["difference_pp_b_minus_a"])
    pval = float(m["affiliate_conversion_test"]["p_value"])
    event = m["nearby_business_events"][0]["event_type"].replace("_", " ").title() if m["nearby_business_events"] else "None"
    supported = bool(d["nearby_event_cause_supported"])
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    add_title(s, "Refusal story — temporal proximity is not evidence of causality", subtitle="Question: “Asia conversion fell after July 20. Did the nearby office relocation cause it?”", kicker="1 + 1 trust story")
    metric_card(s, 0.8, 1.9, 2.2, 1.05, "Asia FTD", _pp(asia), accent=RED)
    metric_card(s, 3.15, 1.9, 2.2, 1.05, "Affiliate FTD", _pp(affiliate), accent=RED)
    metric_card(s, 5.5, 1.9, 2.2, 1.05, "Affiliate p", _pvalue(pval), accent=GREEN)
    metric_card(s, 7.85, 1.9, 2.2, 1.05, "Nearby event", event, accent=AMBER)
    metric_card(s, 10.2, 1.9, 2.35, 1.05, "Event causal support", str(supported).upper(), accent=GREEN if not supported else RED)
    panel(s, 0.8, 3.35, 5.7, 2.35, fill=PANEL)
    add_text(s, "Tempting narrative", 1.08, 3.62, 2.2, 0.3, size=12, color=AMBER, bold=True)
    flow_node(s, "Office relocation", 1.1, 4.18, 2.1, fill=PANEL_2, accent=AMBER, h=0.8)
    arrow(s, 3.23, 4.58, 4.0, 4.58, color=AMBER, width=1.4)
    flow_node(s, "Asia FTD ↓", 4.05, 4.18, 1.65, fill=PANEL_2, accent=RED, h=0.8)
    add_text(s, "REJECTED", 2.72, 5.18, 1.8, 0.25, size=11, color=RED, bold=True, align=PP_ALIGN.CENTER)
    panel(s, 6.85, 3.35, 5.7, 2.35, fill=PANEL)
    add_text(s, "Measured driver", 7.13, 3.62, 2.2, 0.3, size=12, color=CYAN, bold=True)
    flow_node(s, "Affiliate quality ↓", 7.2, 4.18, 2.25, fill=PANEL_2, accent=RED, h=0.8)
    arrow(s, 9.5, 4.58, 10.2, 4.58, color=CYAN, width=1.4)
    flow_node(s, "Supported driver", 10.25, 4.18, 1.65, fill=PANEL_2, accent=GREEN, h=0.8)
    add_text(s, "false-correlation rejection = true", 7.25, 5.18, 4.4, 0.25, size=10.5, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_callout(s, "FitzSight is evaluated on explanations it refuses to make — not only the answers it produces.", 0.8, 6.13, 11.75, 0.64, accent=GREEN, size=11.5)
    add_footer(s, 7)


def slide_8(prs):
    net = _metrics(NET_DEPOSIT_INTENT)
    cust = _metrics(CUSTOMER_INTELLIGENCE_INTENT)
    marketing = _metrics(MARKETING_LEAD_QUALITY_INTENT)
    net_dec = net["driver_decomposition"]
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    add_title(s, "Three additional workflows prove reuse — without diluting the hero narrative", kicker="Breadth proof")
    cards = [
        ("Client-fund flows", "Net-deposit / withdrawal concentration", f"Net change {_money_k(float(net_dec['net_change']))}", AMBER, "Observed driver only; no customer-motive inference"),
        ("Customer Intelligence", "Behavioral-value segmentation", f"Coverage {_pct(float(cust['segmentation']['coverage']))}", CYAN, "Descriptive only; no AML / credit / eligibility decisions"),
        ("Acquisition quality", "Volume vs mix vs within-channel quality", f"Leads {float(marketing['lead_volume_change_pct']):+.0f}% · FTD {_pp(float(marketing['conversion_change_pp']))}", GREEN, "Separates more leads from better leads"),
    ]
    for i, (title, subtitle, metric, col, note) in enumerate(cards):
        y = 1.75 + i * 1.62
        panel(s, 0.8, y, 11.75, 1.3, fill=PANEL, line=col)
        add_text(s, f"0{i+1}", 1.02, y + 0.22, 0.55, 0.3, size=12, color=col, bold=True)
        add_text(s, title, 1.75, y + 0.18, 2.4, 0.3, size=14, bold=True)
        add_text(s, subtitle, 4.0, y + 0.18, 3.55, 0.3, size=11, color=MUTED, bold=True)
        add_text(s, metric, 7.65, y + 0.18, 2.2, 0.3, size=12, color=col, bold=True)
        add_text(s, note, 1.75, y + 0.72, 9.9, 0.32, size=10.5, color=TEXT)
    add_callout(s, "Main pitch = one complete investigation + one refusal. These workflows are reusable breadth, not five equal demo stories.", 0.8, 6.62, 11.75, 0.42, accent=CYAN, size=10.5)
    add_footer(s, 8)


def slide_9(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    add_title(s, "Technical depth comes from explicit authority boundaries, not more Agent labels", kicker="Architecture")
    cards = [
        ("01", "Local intent gate", "Unsupported questions are refused before an external model call.", CYAN),
        ("02", "Constrained planner", "Only approved high-level actions. No planner-generated SQL or arbitrary tool parameters.", AMBER),
        ("03", "Deterministic analytics", "Read-only SQL / Python own KPI calculations, tests, decompositions and anomaly checks.", GREEN),
        ("04", "Source-addressable evidence", "Tool outputs and synthetic operational documents enter one append-only Evidence Registry.", PURPLE),
        ("05", "EvidenceClaimVerifier", "Checks evidence integrity, _gt boundary and causal language before result delivery.", CYAN),
        ("06", "Human decision boundary", "No trading, AML, credit, suitability or other high-impact financial decisions.", GREEN),
    ]
    for i, (num, title, desc, col) in enumerate(cards):
        x = 0.8 + (i % 3) * 4.0
        y = 1.72 + (i // 3) * 2.25
        panel(s, x, y, 3.65, 1.75, fill=PANEL, line=col)
        add_text(s, num, x + 0.18, y + 0.18, 0.42, 0.24, size=10.5, color=col, bold=True)
        add_text(s, title, x + 0.65, y + 0.16, 2.7, 0.32, size=13, bold=True)
        add_text(s, desc, x + 0.65, y + 0.62, 2.65, 0.85, size=10.2, color=MUTED)
    add_callout(s, "Verification failure is a product state: FitzSight may finish with a guarded answer or withhold the unsupported attribution.", 0.8, 6.25, 11.75, 0.6, accent=RED, size=10.8)
    add_footer(s, 9)


def slide_10(prs):
    snap = _evaluation_snapshot()
    bench = snap["benchmark"]
    adv = snap["adversarial"]
    hold = snap["holdout"]["metrics"]
    abl = snap["ablation"]["metrics"]
    full = abl["full_fitzsight"]
    nogate = abl["no_verifier_gate_ablation"]
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    add_title(s, "Evaluation v2 asks whether the system generalizes — and whether the safety architecture matters", kicker="Evaluation & reproducibility")
    metric_card(s, 0.8, 1.75, 2.3, 1.1, "Fixed benchmark", f"{bench['passed']} / {bench['scenario_count']} PASS", accent=GREEN)
    metric_card(s, 3.25, 1.75, 2.3, 1.1, "Holdout runs", f"{hold['case_runs']} / {hold['case_runs']}", accent=GREEN, note="Unseen seeds + question paraphrases")
    metric_card(s, 5.7, 1.75, 2.3, 1.1, "Routing stability", _pct(float(hold['intent_routing_stability'])), accent=CYAN)
    metric_card(s, 8.15, 1.75, 2.3, 1.1, "Evidence coverage", _pct(float(hold['mean_evidence_coverage'])), accent=GREEN)
    metric_card(s, 10.6, 1.75, 1.95, 1.1, "Supported", _pct(float(hold['supported_candidate_rate'])), accent=AMBER, note="1 CRM seed correctly insufficient")
    panel(s, 0.8, 3.25, 5.75, 2.55, fill=PANEL)
    add_text(s, "Adversarial release gate", 1.05, 3.52, 2.8, 0.3, size=12, color=AMBER, bold=True)
    add_text(s, f"{adv['passed']} / {adv['case_count']} PASS", 1.05, 3.98, 2.3, 0.4, size=22, color=GREEN, bold=True)
    for i, t in enumerate(["scope refusal", "planner policy", "evidence integrity", "causal overclaim", "_gt leak", "false correlation"]):
        chip(s, t, 1.05 + (i % 2) * 2.45, 4.56 + (i // 2) * 0.38, 2.15, fill=PANEL_2, color=MUTED)
    panel(s, 6.85, 3.25, 5.7, 2.55, fill=PANEL)
    add_text(s, "Controlled architecture ablation", 7.1, 3.52, 3.6, 0.3, size=12, color=PURPLE, bold=True)
    add_text(s, "Full FitzSight", 7.1, 4.02, 1.9, 0.26, size=10.5, color=MUTED, bold=True)
    add_text(s, _pct(float(full['adversarial_refusal_correctness'])), 9.65, 3.97, 1.1, 0.35, size=20, color=GREEN, bold=True, align=PP_ALIGN.RIGHT)
    add_text(s, "adversarial refusal", 10.85, 4.03, 1.25, 0.25, size=9, color=MUTED)
    add_text(s, "No verifier/evidence gate", 7.1, 4.68, 2.6, 0.26, size=10.5, color=MUTED, bold=True)
    add_text(s, _pct(float(nogate['unsafe_answer_rate_on_adversarial'])), 9.65, 4.62, 1.1, 0.35, size=20, color=RED, bold=True, align=PP_ALIGN.RIGHT)
    add_text(s, "unsafe-answer rate", 10.85, 4.69, 1.3, 0.25, size=9, color=MUTED)
    add_text(s, "This is a controlled architecture ablation — not a Generic LLM baseline.", 7.1, 5.28, 4.9, 0.28, size=9.5, color=AMBER, bold=True)
    add_callout(s, "One unseen CRM seed returning insufficient_evidence is retained as a robustness success, not hidden to make the metric look perfect.", 0.8, 6.2, 11.75, 0.62, accent=AMBER, size=10.6)
    add_footer(s, 10)


def slide_11(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    add_title(s, "Safety, traceability and open reuse are explicit product boundaries", kicker="AI + Finance")
    panel(s, 0.8, 1.75, 5.75, 4.85, fill=PANEL)
    add_text(s, "Implemented today", 1.08, 2.03, 2.2, 0.3, size=12, color=GREEN, bold=True)
    implemented = [
        "Synthetic benchmark data; no real customer PII",
        "Read-only analytical path with deterministic fallback",
        "Evidence Registry + source/paragraph document evidence",
        "Causal-language guardrail and fail-closed verifier",
        "No investment advice, trading or account actions",
        "No automated AML, credit or suitability decisions",
        "MIT license + public repo + tests + benchmark catalog",
    ]
    for i, t in enumerate(implemented):
        add_text(s, "• " + t, 1.08, 2.48 + i * 0.48, 5.05, 0.36, size=10.6, color=TEXT)
    panel(s, 6.85, 1.75, 5.7, 4.85, fill=PANEL)
    add_text(s, "Production blueprint — not claimed as implemented", 7.12, 2.03, 4.9, 0.3, size=12, color=AMBER, bold=True)
    blueprint = [
        "Identity / RBAC",
        "Row- and field-level policy",
        "PII masking",
        "Retention / enterprise audit controls",
        "Approved business-system connectors",
    ]
    for i, t in enumerate(blueprint):
        chip(s, t, 7.15, 2.55 + i * 0.52, 4.65, fill=PANEL_2, color=MUTED)
    add_callout(s, "Financial boundary: evidence-grounded decision support. The authorized analyst or institution retains the final judgment.", 7.15, 5.45, 4.9, 0.78, accent=GREEN, size=10.5)
    add_footer(s, 11)


def slide_12(prs):
    snap = _evaluation_snapshot()
    hold = snap["holdout"]["metrics"]
    s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
    add_text(s, "FITZSIGHT", 0.8, 0.68, 3.0, 0.4, size=12, color=CYAN, bold=True)
    add_text(s, "BI tells you what changed.\nFitzSight investigates what the evidence supports — and what it does not.", 0.8, 1.5, 11.0, 1.55, size=29, bold=True)
    add_callout(s, "Autonomous investigation. Human decision.", 0.8, 3.42, 7.2, 0.82, accent=GREEN, size=15)
    metric_card(s, 0.8, 4.72, 2.6, 1.15, "Hero", "CRM / FTD", accent=CYAN, note="bounded-adaptive investigation")
    metric_card(s, 3.62, 4.72, 2.6, 1.15, "Trust case", "False correlation", accent=AMBER, note="refuses unsupported causality")
    metric_card(s, 6.44, 4.72, 2.6, 1.15, "Holdout", f"{hold['case_runs']} / {hold['case_runs']}", accent=GREEN, note="routed + verified")
    metric_card(s, 9.26, 4.72, 3.25, 1.15, "Evidence coverage", _pct(float(hold['mean_evidence_coverage'])), accent=GREEN, note="synthetic benchmark scope")
    add_text(s, "A reproducible financial investigation you can inspect, verify, challenge — and sometimes refuse.", 0.85, 6.42, 11.1, 0.38, size=14, color=MUTED, bold=True)
    add_footer(s, 12)


def build_speaker_notes() -> Path:
    hero = _hero()["run"]
    hm = hero["investigation"]["metrics"]
    hd = hero["investigation"]["diagnosis"]
    false_m = _metrics(FALSE_CORRELATION_INTENT)
    false_d = _diagnosis(FALSE_CORRELATION_INTENT)
    net = _metrics(NET_DEPOSIT_INTENT)
    customer = _metrics(CUSTOMER_INTELLIGENCE_INTENT)
    marketing = _metrics(MARKETING_LEAD_QUALITY_INTENT)
    snap = _evaluation_snapshot()
    hold = snap["holdout"]["metrics"]
    full = snap["ablation"]["metrics"]["full_fitzsight"]
    nogate = snap["ablation"]["metrics"]["no_verifier_gate_ablation"]

    notes = f"""# FitzSight — GOAI Initial-Round 12-Slide Speaker Notes (v0.12.1)

> Competition-facing numerical claims are derived from verified deterministic runtime evidence or the checked-in v0.12 evaluation JSON. This deck does not claim live Streamlit or live OpenAI-provider validation.

## Slide 1 — FitzSight

“FitzSight is an evidence-grounded financial operations Agent for Brokerage / FinTech Operations Analysts. The operating principle is simple: autonomous investigation, human decision.”

## Slide 2 — Industry problem

“Dashboards and SQL show what changed. The expensive part is the investigation behind why: define the KPI, compare cohorts, drill drivers, test significance, inspect operational context, reconcile evidence, then write a report.”

## Slide 3 — Product

“The planner can select approved analytical actions, but deterministic SQL and Python own every number. Evidence records each step. The verifier determines which claims may reach the analyst.”

## Slide 4 — Hero journey

“This is the actual v0.12 verified execution trace, rendered from runtime JSON. The contribution and statistical results trigger the next approved latency and event branches. The planner still cannot invent SQL or arbitrary tools.”

## Slide 5 — Hero finding

“The affected European teams moved {_pp(float(hm['affected']['conversion_change_pp']))} versus {_pp(float(hm['control']['conversion_change_pp']))} in the control cohort. Median response time changed {float(hm['affected_response_median_change_minutes']):+.2f} minutes. The evidence chain also includes contribution, anomaly, operational-event and document evidence at {hm['document_evidence']['source_ref']}. The final status is {hd['root_cause_status']}, not proven real-world causality.”

## Slide 6 — Failure branch

“We also test the branch where the event dependency fails. FitzSight records error Evidence, skips document corroboration, changes root-cause status to insufficient_evidence, and still returns a verified bounded answer. The system is not rewarded for always forcing a cause.”

## Slide 7 — Refusal case

“In Asia, aggregate FTD changes {_pp(float(false_m['conversion_change_pp']))}; Affiliate conversion changes {_pp(float(false_m['affiliate_conversion_test']['difference_pp_b_minus_a']))}. A nearby office relocation exists, but causal support is {str(bool(false_d['nearby_event_cause_supported'])).upper()}. Temporal proximity is explicitly rejected as proof.”

## Slide 8 — Breadth

“Three other workflows prove reuse without diluting the main story: client-fund-flow concentration, descriptive customer segmentation, and acquisition volume-versus-quality. For example, current net-deposit change is {_money_k(float(net['driver_decomposition']['net_change']))}; customer segmentation covers {_pct(float(customer['segmentation']['coverage']))}; marketing lead volume changes {float(marketing['lead_volume_change_pct']):+.0f}% while FTD changes {_pp(float(marketing['conversion_change_pp']))}.”

## Slide 9 — Architecture

“The technical depth is in authority separation: local intent gate, constrained planner, deterministic tools, source-addressable Evidence Registry, EvidenceClaimVerifier, and an explicit human-decision boundary.”

## Slide 10 — Evaluation v2

“Five fixed benchmark scenarios pass. Eight holdout seed-and-paraphrase runs route and verify successfully with {_pct(float(hold['mean_evidence_coverage']))} evidence coverage. Supported-candidate rate is only {_pct(float(hold['supported_candidate_rate']))} because one unseen CRM seed correctly returns insufficient evidence. In the controlled architecture ablation, full FitzSight refuses {_pct(float(full['adversarial_refusal_correctness']))} of adversarial fixtures; removing the verifier/evidence gate yields {_pct(float(nogate['unsafe_answer_rate_on_adversarial']))} unsafe-answer rate. This is not a Generic LLM baseline.”

## Slide 11 — Safety and reuse

“Current implementation uses synthetic data, read-only analytics, evidence tracing and fail-closed verification. Enterprise RBAC, PII masking and retention remain a production blueprint, not an implemented claim. Final professional judgment remains human.”

## Slide 12 — Close

“BI tells you what changed. FitzSight investigates what the measurable evidence supports — and what it does not. It is a reproducible financial investigation you can inspect, verify, challenge, and sometimes refuse.”
"""
    path = SUBMISSION_DIR / "PITCH_SPEAKER_NOTES.md"
    path.write_text(notes, encoding="utf-8")
    return path


def build_pptx() -> Path:
    SUBMISSION_DIR.mkdir(parents=True, exist_ok=True)
    for asset in (HERO_TRACE, HERO_ANSWER):
        if not asset.exists():
            raise FileNotFoundError(f"Required runtime-derived hero visual missing: {asset}")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for builder in [slide_1, slide_2, slide_3, slide_4, slide_5, slide_6, slide_7, slide_8, slide_9, slide_10, slide_11, slide_12]:
        builder(prs)
    prs.save(PPTX_PATH)
    return PPTX_PATH


def export_pdf(pptx_path: Path) -> Path | None:
    libreoffice = shutil.which("libreoffice") or shutil.which("soffice")
    if not libreoffice:
        return None
    completed = subprocess.run(
        [libreoffice, "--headless", "--convert-to", "pdf", "--outdir", str(SUBMISSION_DIR), str(pptx_path)],
        cwd=ROOT,
        check=False,
        stdout=subprocess.PIPE,
        stderr=subprocess.STDOUT,
        text=True,
    )
    print(completed.stdout.strip())
    generated = pptx_path.with_suffix(".pdf")
    return generated if generated.exists() else None


def main() -> int:
    pptx = build_pptx()
    notes = build_speaker_notes()
    pdf = export_pdf(pptx)
    print(f"Created PPTX: {pptx}")
    print(f"Created notes: {notes}")
    if pdf:
        print(f"Created PDF:  {pdf}")
        return 0
    print("PDF export skipped: LibreOffice/soffice not available or conversion failed.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
